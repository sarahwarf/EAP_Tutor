"""Extract plain text from PDF, PPTX, DOCX, and TXT files."""
import io


def extract(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """
    Returns (text, file_type) from raw file bytes.
    file_type is one of: pdf, pptx, docx, txt
    """
    name = filename.lower()

    if name.endswith(".pdf"):
        return _from_pdf(file_bytes), "pdf"
    elif name.endswith(".pptx"):
        return _from_pptx(file_bytes), "pptx"
    elif name.endswith(".docx"):
        return _from_docx(file_bytes), "docx"
    else:
        return file_bytes.decode("utf-8", errors="ignore"), "txt"


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages).strip()


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()


def _from_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs).strip()
